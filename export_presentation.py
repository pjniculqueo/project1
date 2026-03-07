import os
from pptx import Presentation
from pptx.util import Inches, Pt


def crear_presentacion(img_dir="resultados_simulacion", output_file="Presentacion_Ministro.pptx"):
    prs = Presentation()

    # Configuramos el tamaño de la diapositiva a pantalla panorámica (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Diapositiva 1: Título ---
    slide_layout = prs.slide_layouts[0]  # Diseño de título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Impacto Cuantitativo de los Fondos Generacionales y la PGU"
    subtitle.text = "Evidencia Causal y Optimización de Política Pública\nMinisterio de Trabajo y Previsión Social"

    # Función auxiliar para agregar diapositivas de contenido + imágenes
    def add_content_slide(title_text, bullet_points, image_files):
        layout = prs.slide_layouts[5]  # Diseño de solo título
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title_text

        # Agregar cuadro de texto en la izquierda
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, point in enumerate(bullet_points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(18)
            p.level = 0

        # Agregar imágenes en la derecha
        img_left = Inches(5.8)
        img_top = Inches(1.5)
        img_width = Inches(7)

        for img_name in image_files:
            img_path = os.path.join(img_dir, img_name)
            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
                img_top += Inches(3)  # Bajar la siguiente imagen si hay más de una
            else:
                print(f"⚠️ Advertencia: No se encontró la imagen {img_path}")

    print("Generando diapositivas...")

    # --- Diapositiva 2 ---
    add_content_slide(
        "El Cambio de Paradigma en las Pensiones",
        [
            "• Chile transita de un modelo de elección (Multifondos) a uno determinista (Fondos Generacionales).",
            "• El Glidepath obliga a migrar a renta fija en los últimos 15 años laborales.",
            "• El modelo IA (GAN) demuestra que estrategias activas mantienen riesgo para aprovechar el interés compuesto.",
            "• El nuevo esquema institucionaliza una estrategia rígida que ignora ciclos económicos."
        ],
        ["plot2A_ponderaciones_gan.png"]
    )

    # --- Diapositiva 3 ---
    add_content_slide(
        "El Costo Oculto de la 'Seguridad Obligatoria'",
        [
            "• Logro de la Ley: Recorta la cola de pérdidas (Mejora el Value at Risk).",
            "• Costo de Oportunidad: Destruye el crecimiento exponencial en el periodo de mayor capital.",
            "• Conclusión: El Fondo Generacional es un seguro contra crisis, pero se paga con una severa caída en la pensión promedio."
        ],
        ["plot3_rendimiento.png"]
    )

    # --- Diapositiva 4 ---
    add_content_slide(
        "La PGU como el Verdadero Seguro contra Ruina",
        [
            "• La PGU ya trunca matemáticamente el riesgo de cola para la mayoría de la población.",
            "• Forzar a sectores vulnerables a refugiarse en renta fija genera sobreprotección redundante.",
            "• Al bajar su riesgo, bajan sus retornos, anclándolos permanentemente a la dependencia del subsidio estatal máximo."
        ],
        ["plot5_impacto_pgu.png"]
    )

    # --- Diapositiva 5 ---
    add_content_slide(
        "Sensibilidad Etaria y la Tensión del Retiro",
        [
            "• Jubilar a los 65 años en el nuevo sistema TDF reduce dramáticamente la pensión frente a una estrategia optimizada.",
            "• Para mantener tasas de reemplazo dignas, el diseño forzará de facto a las personas a postergar su jubilación hacia los 70 años."
        ],
        ["plot6_edades_jubilacion.png"]
    )

    # --- Diapositiva 6 ---
    add_content_slide(
        "Propuesta: Riesgo Diferenciado Optimizado",
        [
            "• Abandonar el Glidepath estándar y transicionar a una asignación condicionada al ingreso.",
            "• Bajo Ingreso: Mantener riesgo al 85%. Si falla, el Estado paga PGU (riesgo neutral). Si sube, sale del subsidio.",
            "• Alto Ingreso: Reducir riesgo a 0%. Evita que crisis de última hora los conviertan en carga fiscal.",
            "• Resultado A/B Testing: Ahorro para el fisco y mayor pensión para los afiliados."
        ],
        ["plot7_optimizacion_riesgo_pgu.png"]
    )

    prs.save(output_file)
    print(f"✅ ¡Presentación exportada con éxito! Archivo: {os.path.abspath(output_file)}")


if __name__ == "__main__":
    crear_presentacion()